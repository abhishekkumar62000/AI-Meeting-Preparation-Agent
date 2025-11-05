import io
import base64
import json
import os
from datetime import datetime
from pathlib import Path
from typing import Dict, List

import streamlit as st
from crewai import Agent, Task, Crew, LLM
from crewai.process import Process
from crewai_tools import SerperDevTool
from pypdf import PdfReader


DATA_DIR = Path(__file__).parent / "data"
HISTORY_FILE = DATA_DIR / "meeting_history.json"
MAX_DOCUMENT_CHARS = 6000


def _read_history_file() -> List[Dict[str, str]]:
    if not HISTORY_FILE.exists():
        return []
    try:
        return json.loads(HISTORY_FILE.read_text(encoding="utf-8"))
    except Exception:
        return []


def _result_to_markdown(result) -> str:
    """Best-effort conversion of CrewAI results to displayable markdown/text.

    Handles CrewOutput objects by checking common attributes and falling back to str().
    Ensures outputs are JSON-serializable when persisted in history.
    """
    # Direct string
    if isinstance(result, str):
        return result
    # Common CrewAI/CrewOutput fields that may contain text
    for attr in ("raw", "final_output", "result", "output"):
        try:
            val = getattr(result, attr, None)
        except Exception:
            val = None
        if isinstance(val, str) and val.strip():
            return val
    # Dict-like export
    if hasattr(result, "to_dict"):
        try:
            return json.dumps(result.to_dict(), indent=2, ensure_ascii=False)
        except Exception:
            pass
    # Fallback
    try:
        return str(result)
    except Exception:
        return ""


def _brief_to_pptx(markdown_text: str, title: str = "Meeting Brief") -> bytes:
    """Convert a markdown brief into a simple PPTX deck and return bytes.

    Creates a title slide and subsequent slides for top-level headings and bullets.
    If python-pptx is not installed, returns empty bytes.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:
        # Dependency missing; return empty and let caller warn
        return b""

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Auto-generated from executive brief"

    # Naive markdown parsing for headings and bullets
    lines = markdown_text.splitlines()
    current_slide = None
    bullet_layout = prs.slide_layouts[1]

    def add_slide_with_title(heading: str):
        nonlocal current_slide
        current_slide = prs.slides.add_slide(bullet_layout)
        current_slide.shapes.title.text = heading.strip("# ").strip()

    def add_bullet(text: str):
        if not current_slide:
            add_slide_with_title("Details")
        body = current_slide.shapes.placeholders[1].text_frame
        if len(body.paragraphs) == 1 and not body.paragraphs[0].text:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = text.strip("- *\t ")
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)

    for line in lines:
        if line.strip().startswith("# ") or line.strip().startswith("## "):
            add_slide_with_title(line)
        elif line.strip().startswith(("- ", "* ")):
            add_bullet(line)
        elif line.strip():
            # treat as a bullet for simplicity
            add_bullet(line)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _write_history_file(entries: List[Dict[str, str]]) -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    HISTORY_FILE.write_text(json.dumps(entries, indent=2), encoding="utf-8")


def _truncate_text(value: str, max_chars: int = MAX_DOCUMENT_CHARS) -> str:
    if len(value) <= max_chars:
        return value
    return value[:max_chars] + "\n\n...[truncated]..."


def _extract_supporting_documents(uploaded_files) -> List[Dict[str, str]]:
    documents: List[Dict[str, str]] = []
    for uploaded in uploaded_files or []:
        file_bytes = uploaded.getvalue()
        if not file_bytes:
            continue
        name = uploaded.name
        text = ""
        if name.lower().endswith(".pdf"):
            try:
                reader = PdfReader(io.BytesIO(file_bytes))
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
            except Exception as exc:
                st.warning(f"Could not read {name}: {exc}")
                continue
        else:
            try:
                text = file_bytes.decode("utf-8", errors="ignore")
            except Exception as exc:
                st.warning(f"Could not decode {name}: {exc}")
                continue
        text = text.strip()
        if text:
            documents.append({"name": name, "content": text})
    return documents


def _build_document_digest(documents: List[Dict[str, str]]) -> str:
    if not documents:
        return "No additional supporting documents were provided."
    sections = []
    for doc in documents:
        sections.append(
            f"Document: {doc['name']}\nContent Preview:\n{_truncate_text(doc['content'], st.session_state.get('truncate_chars', MAX_DOCUMENT_CHARS))}"
        )
    return "\n\n".join(sections)


def _save_meeting_to_history(entry: Dict[str, str]) -> None:
    entries = st.session_state.get("meeting_history", [])
    entries.insert(0, entry)
    entries = entries[:20]
    st.session_state["meeting_history"] = entries
    _write_history_file(entries)


def _format_history_option(entry: Dict[str, str]) -> str:
    company = entry.get("company", "Unknown")
    objective = entry.get("objective", "No objective")
    timestamp = entry.get("timestamp", "No timestamp")
    if isinstance(timestamp, str) and "T" in timestamp:
        timestamp = timestamp.replace("T", " ")
    return f"{company} – {objective} ({timestamp})"


def _clear_history() -> None:
    try:
        if HISTORY_FILE.exists():
            HISTORY_FILE.unlink()
    except Exception:
        pass
    st.session_state["meeting_history"] = []
    st.session_state["history_view"] = None


# Streamlit app setup
st.set_page_config(page_title="AI Meeting Agent 📝", layout="wide")


def _inject_theme_css(preset: str = "Neon Night") -> None:
    """Inject custom CSS for multiple colorful themes and subtle animations.

    Presets: 'Neon Night', 'Emerald Dark', 'Sunset Dark', 'Light Minimal'
    """
    p = preset.lower()
    if p == "emerald dark":
        primary, accent = "#34d399", "#22d3ee"
        bg, bg2, text, subtle = "#0a0f0f", "#0f1717", "#e7f6f2", "#7aa3a3"
    elif p == "sunset dark":
        primary, accent = "#fb7185", "#f59e0b"
        bg, bg2, text, subtle = "#0f0b10", "#17101a", "#f5e9f1", "#b08aa6"
    elif p == "light minimal":
        primary, accent = "#6366f1", "#06b6d4"
        bg, bg2, text, subtle = "#f7f8fc", "#ffffff", "#0b0f19", "#475569"
    else:  # Neon Night (default)
        primary, accent = "#9b87f5", "#22d3ee"
        bg, bg2, text, subtle = "#0b0f19", "#121829", "#E6EAF2", "#94a3b8"

        st.markdown(
                f"""
                <style>
                :root {{
                    --primary: {primary};
                    --accent: {accent};
                    --bg: {bg};
                    --bg2: {bg2};
                    --text: {text};
                    --subtle: {subtle};
                }}

                /* App background with soft gradient */
                .stApp {{
                    background: radial-gradient(1200px 600px at 10% -10%, rgba(34,211,238,0.08), transparent 60%),
                                            radial-gradient(1000px 500px at 90% 0%, rgba(155,135,245,0.12), transparent 50%),
                                            var(--bg);
                    color: var(--text);
                }}

                        /* Hero banner */
                        .hero-wrap {{
                            padding: 12px 18px; border-radius: 14px; border: 1px solid rgba(148,163,184,0.18);
                            background: linear-gradient(90deg, rgba(155,135,245,0.12), rgba(34,211,238,0.12));
                            display:flex; flex-direction:column; gap:6px; margin-bottom: 8px;
                            animation: fadeIn .35s ease;
                        }}
                        .hero-title {{
                            font-size: 28px; font-weight: 800; line-height: 1.15;
                            background: linear-gradient(90deg, var(--primary), var(--accent));
                            -webkit-background-clip: text; background-clip: text; color: transparent;
                            background-size: 200% 100%;
                            animation: gradientMove 8s ease infinite, titleGlow 3s ease-in-out infinite alternate;
                        }}
                        .hero-subtitle {{
                            opacity:.0; font-size: 14px; color: var(--text);
                            text-shadow: 0 0 0 rgba(34,211,238,0);
                            animation: subtitleIn .6s ease forwards .2s, subtitlePulse 6s ease-in-out infinite 1.2s;
                        }}
                        .hero-underline {{
                            height: 2px; width: 120px; border-radius: 999px;
                            background: linear-gradient(90deg, var(--primary), var(--accent));
                            background-size: 200% 100%;
                            animation: gradientMove 8s ease infinite;
                        }}

                        /* Title gradient text */
                        h1 {{
                    background: linear-gradient(90deg, var(--primary), var(--accent));
                    -webkit-background-clip: text; background-clip: text; color: transparent;
                    animation: gradientMove 8s ease infinite;
                    background-size: 200% 100%;
                }}

                /* Cards (expanders, sidebar blocks) */
                .st-expander, .st-emotion-cache-1r6slb0, .st-emotion-cache-16idsys {{
                    background: linear-gradient(180deg, rgba(255,255,255,0.02), rgba(255,255,255,0.00));
                    border: 1px solid rgba(148,163,184,0.18);
                    border-radius: 14px;
                }}
                .stExpanderHeader, .st-emotion-cache-1r6slb0, .st-emotion-cache-1avcm0n {{
                    color: var(--text);
                }}

                        /* Buttons: primary glow + hover lift */
                .stButton>button, .stDownloadButton>button {{
                    background: linear-gradient(90deg, var(--primary), var(--accent));
                    border: none; color: #0b0f19; font-weight: 700;
                    border-radius: 10px; padding: 0.6rem 1rem;
                    box-shadow: 0 8px 20px rgba(34,211,238,0.18);
                    transition: transform .12s ease, box-shadow .2s ease, filter .2s ease;
                }}
                .stButton>button:hover, .stDownloadButton>button:hover {{
                    transform: translateY(-1px);
                    box-shadow: 0 10px 24px rgba(155,135,245,0.25);
                    filter: brightness(1.03);
                }}
                .stButton>button:active {{ transform: translateY(0); }}

                        /* Inputs */
                        .stTextInput>div>div>input, textarea, .stNumberInput input {{
                    background: var(--bg2); color: var(--text);
                    border-radius: 10px; border: 1px solid rgba(148,163,184,0.25);
                }}
                        .stSelectbox div[data-baseweb="select"]>div {{
                            background: var(--bg2);
                            color: var(--text);
                            border: 1px solid rgba(148,163,184,0.25);
                            border-radius: 10px;
                        }}
                        .stMultiSelect div[data-baseweb="select"]>div {{
                            background: var(--bg2);
                            color: var(--text);
                            border: 1px solid rgba(148,163,184,0.25);
                            border-radius: 10px;
                        }}
                        /* Slider */
                        div[data-testid="stSlider"] .stSlider>div>div>div>div {{
                            background: linear-gradient(90deg, var(--primary), var(--accent));
                        }}
                        /* Checkboxes/Radio */
                        .stCheckbox>label>div>span, .stRadio>label {{ color: var(--text); }}
                        /* File Uploader */
                        section[data-testid="stFileUploader"] div[role="button"] {{
                            background: var(--bg2);
                            border: 1px dashed rgba(148,163,184,0.35);
                        }}
                        /* Tables */
                        .stTable tbody tr:hover {{ background: rgba(148,163,184,0.08); }}
                        .stTable th {{ color: var(--text); }}

                /* Sidebar */
                section[data-testid="stSidebar"]>div {{
                    background: linear-gradient(180deg, rgba(155,135,245,0.06), rgba(34,211,238,0.03));
                    border-right: 1px solid rgba(148,163,184,0.18);
                }}

                        /* Container fade-in */
                        .block-container {{
                            animation: fadeIn .35s ease;
                        }}

                                /* Sidebar logo */
                                .sidebar-logo-wrap {{
                                    display:flex; align-items:center; justify-content:center;
                                    padding: 12px 6px 6px 6px; margin-bottom: 10px;
                                }}
                                .sidebar-logo {{
                                    max-width: 160px; width: 100%; height:auto; border-radius: 12px;
                                    box-shadow: 0 10px 24px rgba(0,0,0,0.25);
                                    animation: floatY 6s ease-in-out infinite;
                                    transition: transform .18s ease, filter .25s ease;
                                }}
                                .sidebar-logo:hover {{
                                    transform: translateY(-2px) scale(1.02);
                                    filter: drop-shadow(0 6px 18px rgba(34,211,238,0.25));
                                }}

                /* Animated subtle gradient */
                @keyframes gradientMove {{
                    0% {{background-position: 0% 50%;}}
                    50% {{background-position: 100% 50%;}}
                    100% {{background-position: 0% 50%;}}
                }}
                                @keyframes fadeIn {{
                                    from {{ opacity: 0; transform: translateY(4px); }}
                                    to {{ opacity: 1; transform: translateY(0); }}
                                }}
                        @keyframes titleGlow {{
                            from {{ text-shadow: 0 0 6px rgba(155,135,245,0.18); }}
                            to {{ text-shadow: 0 0 14px rgba(34,211,238,0.28); }}
                        }}
                        @keyframes subtitleIn {{
                            from {{ opacity: 0; transform: translateY(4px); }}
                            to {{ opacity: 1; transform: translateY(0); }}
                        }}
                        @keyframes subtitlePulse {{
                            0% {{ text-shadow: 0 0 0 rgba(34,211,238,0); }}
                            50% {{ text-shadow: 0 0 10px rgba(34,211,238,0.20); }}
                            100% {{ text-shadow: 0 0 0 rgba(34,211,238,0); }}
                        }}
                                @keyframes floatY {{
                                    0% {{ transform: translateY(0); }}
                                    50% {{ transform: translateY(-4px); }}
                                    100% {{ transform: translateY(0); }}
                                }}
                </style>
                """,
                unsafe_allow_html=True,
        )


st.markdown(
        """
        <div class="hero-wrap">
            <div class="hero-title">AI Meeting Preparation Agent 📝</div>
            <div class="hero-subtitle">Multi-agent intelligence for real meetings Turn meetings into momentum.</div>
            <div class="hero-underline"></div>
        </div>
        """,
        unsafe_allow_html=True,
)

if "meeting_history" not in st.session_state:
    st.session_state["meeting_history"] = _read_history_file()
if "history_view" not in st.session_state:
    st.session_state["history_view"] = None

# Sidebar for API keys
def _render_sidebar_logo():
    try:
        logo_path = Path(__file__).parent / "Logo.png"
        if logo_path.exists():
            b64 = base64.b64encode(logo_path.read_bytes()).decode("ascii")
            st.sidebar.markdown(
                f"""
                <div class="sidebar-logo-wrap">
                  <img class="sidebar-logo" src="data:image/png;base64,{b64}" alt="App Logo" />
                </div>
                """,
                unsafe_allow_html=True,
            )
        else:
            # Minimal fallback if image missing
            st.sidebar.markdown(
                """
                <div class="sidebar-logo-wrap" title="Logo.png not found">
                  <div style="width:96px;height:96px;border-radius:18px; 
                              background: linear-gradient(135deg, var(--primary), var(--accent));
                              box-shadow: 0 10px 24px rgba(0,0,0,0.25);
                              display:flex;align-items:center;justify-content:center; 
                              font-weight:800;color:#0b0f19;">AI</div>
                </div>
                """,
                unsafe_allow_html=True,
            )
    except Exception:
        pass

_render_sidebar_logo()

st.sidebar.header("API Keys")
openai_api_key = st.sidebar.text_input("OpenAI API Key", type="password")
serper_api_key = st.sidebar.text_input("Serper API Key", type="password")

st.sidebar.subheader("Appearance")
theme_preset = st.sidebar.selectbox(
    "Theme Preset",
    ["Neon Night", "Emerald Dark", "Sunset Dark", "Light Minimal"],
    index=0,
    key="theme_preset",
)
_inject_theme_css(theme_preset)

st.sidebar.subheader("Model Settings")
model_name = st.sidebar.selectbox(
    "LLM model",
    options=[
        "gpt-4.1",
        "gpt-4.1-mini",
        "gpt-4.1-nano",
    ],
    index=2,
)
temperature_setting = st.sidebar.slider("Temperature", 0.0, 1.0, 0.7, 0.05)
st.sidebar.subheader("Document Settings")
st.session_state["truncate_chars"] = st.sidebar.slider(
    "Document preview length (chars)", 1000, 20000, MAX_DOCUMENT_CHARS, 500
)

st.sidebar.subheader("Meeting Library")
history_entries = st.session_state.get("meeting_history", [])
if history_entries:
    selected_idx = st.sidebar.selectbox(
        "Saved briefs",
        options=list(range(len(history_entries))),
        format_func=lambda idx: _format_history_option(history_entries[idx]),
        key="history_selector"
    )
    st.session_state["history_view"] = history_entries[selected_idx]
    if st.sidebar.button("Load into form", use_container_width=True):
        selected = history_entries[selected_idx]
        st.session_state["company_name"] = selected.get("company", "")
        st.session_state["meeting_objective"] = selected.get("objective", "")
        st.session_state["attendees"] = selected.get("attendees", "")
        st.session_state["focus_areas"] = selected.get("focusAreas", "")
else:
    st.sidebar.info("No saved briefs yet. Generate one to get started.")
    st.session_state["history_view"] = None

with st.sidebar.expander("Library actions"):
    confirm_clear = st.checkbox("Confirm clear library", key="confirm_clear_library")
    if st.button("Clear Library", disabled=not confirm_clear):
        _clear_history()
        st.success("Library cleared.")

# Check if required API keys are set
if openai_api_key and serper_api_key:
    os.environ["OPENAI_API_KEY"] = openai_api_key
    os.environ["SERPER_API_KEY"] = serper_api_key

    llm = LLM(model=model_name, temperature=temperature_setting, api_key=openai_api_key)
    search_tool = SerperDevTool()

    company_name = st.text_input("Enter the company name:", key="company_name")
    meeting_objective = st.text_input("Enter the meeting objective:", key="meeting_objective")
    attendees = st.text_area("Enter the attendees and their roles (one per line):", key="attendees")
    meeting_duration = st.number_input(
        "Enter the meeting duration (in minutes):",
        min_value=15,
        max_value=180,
        value=60,
        step=15,
        key="meeting_duration"
    )
    focus_areas = st.text_area("Enter any specific areas of focus or concerns:", key="focus_areas")

    with st.expander("Advanced preparation inputs"):
        uploaded_files = st.file_uploader(
            "Upload supporting documents (PDF, TXT, or MD)",
            type=["pdf", "txt", "md"],
            accept_multiple_files=True
        )
        meeting_notes = st.text_area(
            "Paste relevant historical notes, CRM snippets, or transcript highlights:",
            height=140,
            key="meeting_notes"
        )
        attendee_personas = st.text_area(
            "Describe attendee personas or preferences for personalization:",
            height=100,
            key="attendee_personas"
        )
        rehearsal_focus = st.text_area(
            "List the scenarios or objections you want the rehearsal simulation to cover:",
            height=100,
            key="rehearsal_focus"
        )
        followup_channels = st.text_input(
            "Preferred follow-up channels (e.g. email, Slack, CRM task list):"
        )
        include_live_updates = st.checkbox(
            "Enable real-time web intelligence refresh",
            value=True,
            key="include_live_updates"
        )
        include_regulatory = st.checkbox(
            "Add compliance, localization, and risk insights",
            value=False,
            key="include_regulatory"
        )

        # Nested expanders are not supported by Streamlit. Use a container here
        # to preserve grouping/spacing without creating a nested expander.
        with st.container():
            st.subheader("Paste Invite (auto-extract)")
            invite_text = st.text_area("Paste calendar invite or ICS text", key="paste_invite_text", height=140)
            if st.button("Extract Fields", key="extract_from_invite"):
                # Heuristic parsing for company and attendees
                company_guess = ""
                attendees_guess = []
                duration_guess = None
                try:
                    for line in (invite_text or "").splitlines():
                        ls = line.strip()
                        if not company_guess and ls.lower().startswith(("subject:", "title:")):
                            company_guess = ls.split(":", 1)[-1].strip()
                        if any(tok in ls.lower() for tok in ["attendees", "participants", "with:"]):
                            attendees_guess.append(ls)
                        if any(tok in ls.lower() for tok in ["duration", "minutes", "mins", "min"]):
                            # naive duration capture
                            import re
                            m = re.search(r"(\d{1,3})\s*(min|mins|minutes)", ls.lower())
                            if m:
                                duration_guess = int(m.group(1))
                except Exception:
                    pass

                if company_guess:
                    st.session_state["company_name"] = company_guess[:120]
                if attendees_guess:
                    st.session_state["attendees"] = "\n".join(attendees_guess)[:2000]
                if duration_guess:
                    st.session_state["meeting_duration"] = max(15, min(180, duration_guess))
                st.success("Invite fields extracted into the form.")

    supporting_documents = _extract_supporting_documents(uploaded_files)
    documents_digest = _build_document_digest(supporting_documents)

    if supporting_documents:
        st.info(f"{len(supporting_documents)} supporting document(s) ingested for analysis.")
        with st.expander("Document snapshots"):
            for doc in supporting_documents:
                preview_len = min(st.session_state.get("truncate_chars", MAX_DOCUMENT_CHARS), 800)
                st.markdown(f"**{doc['name']}**\n\n{_truncate_text(doc['content'], preview_len)}")

    context_directives: List[str] = []
    if include_live_updates:
        context_directives.append(
            "Incorporate the most recent news, market movements, and growth signals discovered via live search."
        )
    if include_regulatory:
        context_directives.append(
            "Highlight regulatory, compliance, and localization considerations that could influence the meeting."
        )
    if meeting_notes:
        context_directives.append(
            "Bridge insights with the historical notes provided to emphasize continuity and momentum."
        )

    context_directives_text = "\n".join(f"- {directive}" for directive in context_directives) or "- Focus on actionable intelligence." 

    shared_context = f"""
Supporting documents summary:
{documents_digest}

Historical notes supplied by the team:
{meeting_notes or 'No additional notes were provided.'}
"""

    # Define the agents
    context_analyzer = Agent(
        role="Meeting Context Specialist",
        goal="Analyze and summarize key background information for the meeting",
        backstory="You are an expert at quickly understanding complex business contexts and identifying critical information.",
        verbose=True,
        allow_delegation=False,
        llm=llm,
        tools=[search_tool]
    )

    industry_insights_generator = Agent(
        role="Industry Expert",
        goal="Provide in-depth industry analysis and identify key trends",
        backstory="You are a seasoned industry analyst with a knack for spotting emerging trends and opportunities.",
        verbose=True,
        allow_delegation=False,
        llm=llm,
        tools=[search_tool]
    )

    strategy_formulator = Agent(
        role="Meeting Strategist",
        goal="Develop a tailored meeting strategy and detailed agenda",
        backstory="You are a master meeting planner, known for creating highly effective strategies and agendas.",
        verbose=True,
        allow_delegation=False,
        llm=llm,
    )

    executive_briefing_creator = Agent(
        role="Communication Specialist",
        goal="Synthesize information into concise and impactful briefings",
        backstory="You are an expert communicator, skilled at distilling complex information into clear, actionable insights.",
        verbose=True,
        allow_delegation=False,
        llm=llm,
    )

    rehearsal_coach = Agent(
        role="Executive Rehearsal Coach",
        goal="Simulate the meeting experience and stress-test positioning",
        backstory="You facilitate realistic rehearsals, crafting likely objections and guiding executives on confident responses.",
        verbose=True,
        allow_delegation=False,
        llm=llm,
    )

    follow_up_partner = Agent(
        role="Post-Meeting Activation Partner",
        goal="Translate insights into action items, follow-ups, and enablement assets",
        backstory="You ensure every meeting converts into momentum with crisply defined next steps and communication plans.",
        verbose=True,
        allow_delegation=False,
        llm=llm,
    )

    # Live Transcript (paste) - summarize and extract actions
    with st.expander("Live Transcript (paste)"):
        transcript_text = st.text_area("Paste transcript or large notes here", key="live_transcript_text", height=180)
        col_a, col_b = st.columns(2)
        with col_a:
            extract_now = st.button("Summarize & Extract Actions", key="extract_actions_btn")
        with col_b:
            save_to_library = st.checkbox("Save summary to library", value=False, key="save_transcript_summary")

        if extract_now and transcript_text.strip():
            task = Task(
                description=f"""
                You are assisting during or after a meeting with {company_name}.
                Based on the transcript below, produce:
                - A concise executive summary (5-8 bullets)
                - Decisions made (if any)
                - Action items table: Item, Owner (if detectable), Due (suggested), Success metric
                - Risks and open questions

                Meeting objective: {meeting_objective}
                Attendees: {attendees}
                Focus areas: {focus_areas}

                Transcript:
                {transcript_text}
                """,
                agent=follow_up_partner,
                expected_output="Concise summary with decisions, action items, and risks as markdown."
            )
            temp_crew = Crew(agents=[follow_up_partner], tasks=[task], verbose=False, process=Process.sequential)
            with st.spinner("Extracting action items and summary..."):
                try:
                    tx_result = temp_crew.kickoff()
                except Exception as exc:
                    st.error(f"Error summarizing transcript: {exc}")
                    tx_result = None
            if tx_result:
                tx_text = _result_to_markdown(tx_result)
                st.markdown(tx_text)
                try:
                    st.download_button(
                        label="Download transcript summary",
                        data=tx_text,
                        file_name=f"transcript_summary_{company_name}.md",
                        mime="text/markdown",
                        key="download_tx_summary"
                    )
                except Exception:
                    pass
                if save_to_library:
                    timestamp = datetime.utcnow().isoformat()
                    entry = {
                        "timestamp": timestamp,
                        "company": company_name,
                        "objective": f"{meeting_objective} (Transcript Summary)",
                        "attendees": attendees,
                        "focusAreas": focus_areas,
                        "documents": [],
                        "result": tx_text,
                    }
                    _save_meeting_to_history(entry)
                    st.success("Transcript summary saved to library.")

    # Define the tasks
    context_analysis_task = Task(
        description=f"""
        Analyze the context for the meeting with {company_name}, considering:
        1. The meeting objective: {meeting_objective}
        2. The attendees: {attendees}
        3. The meeting duration: {meeting_duration} minutes
        4. Specific focus areas or concerns: {focus_areas}

        Directives to prioritize:
        {context_directives_text}

        Reference the following shared materials:
        {shared_context}

        Research {company_name} thoroughly, including:
        1. Recent news and press releases (refresh if new headlines are available)
        2. Key products or services
        3. Major competitors and differentiators

        Provide a comprehensive summary of your findings, highlighting the most relevant information for the meeting context.
        Format output using markdown with clear headings and subheadings.
        """,
        agent=context_analyzer,
        expected_output="Detailed meeting context analysis covering company background, latest developments, and insights tied to the objective."
    )

    industry_analysis_task = Task(
        description=f"""
        Based on the context analysis for {company_name} and the meeting objective: {meeting_objective}, provide an in-depth industry analysis:
        1. Identify key trends and developments in the industry
        2. Analyze the competitive landscape and challenger approaches
        3. Highlight potential opportunities and threats for the meeting sponsor
        4. Provide insights on market positioning compared to peers

        Supporting materials to infuse:
        {documents_digest}

        Ensure the analysis is relevant to the meeting objective and attendees' roles.
        Format output using markdown with appropriate headings and subheadings.
        """,
        agent=industry_insights_generator,
        expected_output="Comprehensive industry analysis aligned to the meeting goal, emphasizing opportunities, risks, and differentiation."
    )

    strategy_development_task = Task(
        description=f"""
        Using the context analysis and industry insights, develop a tailored meeting strategy and detailed agenda for the {meeting_duration}-minute meeting with {company_name}. Include:
        1. A time-boxed agenda with clear objectives for each section
        2. Key talking points for each agenda item, connected to both business value and risk mitigation
        3. Suggested speakers or leaders for each section, mapped to {attendees}
        4. Potential discussion topics and questions to drive the conversation
        5. Strategies to address the specific focus areas and concerns: {focus_areas}
        6. Personalization cues leveraging attendee personas: {attendee_personas or 'No persona insights provided.'}

        Ensure the strategy and agenda align with the meeting objective: {meeting_objective}
        Format output using markdown with appropriate headings and subheadings.
        """,
        agent=strategy_formulator,
        expected_output="Detailed meeting strategy and time-boxed agenda mapping objectives to owners, talking points, and success signals."
    )

    executive_brief_task = Task(
        description=f"""
        Synthesize all gathered information into a comprehensive executive brief for the meeting with {company_name}. Create the following components:

        1. A detailed one-page executive summary including:
           - Clear statement of the meeting objective
           - List of key attendees and their roles
           - Critical background points about {company_name} and relevant industry context
           - Top 3-5 strategic goals for the meeting, aligned with the objective
           - Brief overview of the meeting structure and key topics to be covered

        2. An in-depth list of key talking points, each supported by:
           - Relevant data or statistics
           - Specific examples or case studies
           - Connection to the company's current situation or challenges

        3. Anticipate and prepare for potential questions:
           - List likely questions from attendees based on their roles and the meeting objective
           - Craft thoughtful, data-driven responses to each question
           - Include any supporting information or additional context that might be needed

        4. Strategic recommendations and next steps:
           - Provide 3-5 actionable recommendations based on the analysis
           - Outline clear next steps for implementation or follow-up
           - Suggest timelines or deadlines for key actions
           - Identify potential challenges or roadblocks and propose mitigation strategies

        Ensure the brief is comprehensive yet concise, highly actionable, and precisely aligned with the meeting objective: {meeting_objective}. The document should be structured for easy navigation and quick reference during the meeting. Integrate the shared materials below when relevant:
        {shared_context}
        """,
        agent=executive_briefing_creator,
        expected_output="Executive-ready brief including summary, key talking points, risk mitigation, and strategic recommendations."
    )

    rehearsal_simulation_task = Task(
        description=f"""
        Facilitate a rehearsal simulation for the meeting with {company_name}. Deliver:
        1. A scripted dry-run agenda with prompts for each speaker
        2. Persona-driven objections or tough questions informed by {attendees} and the focus areas: {focus_areas}
        3. Suggested high-confidence responses grounded in the research generated so far
        4. Coaching tips on body language, tone, and supporting assets
        5. Scenario branches for unexpected pivots. Prioritize the following rehearsal focus:
           {rehearsal_focus or 'No additional simulation requests provided.'}

        Reference the broader preparation outputs so the rehearsal reflects the planned meeting arc.
        """,
        agent=rehearsal_coach,
        expected_output="Simulation guide featuring persona-based Q&A, objection handling, and coaching cues."
    )

    follow_up_activation_task = Task(
        description=f"""
        Convert the meeting plan into actionable follow-ups. Produce:
        1. A prioritized action item tracker with owners, due dates, and success metrics
        2. A draft follow-up communication tailored to the preferred channel(s): {followup_channels or 'Not specified'}
        3. Recommendations for logging outcomes in CRM or project tools
        4. A checklist for meeting-day capture (notes, decisions, risks, commitments)
        5. Guidance on leveraging any meeting recordings or transcripts post-session

        Integrate the shared preparation context and emphasize how to maintain momentum immediately after the meeting concludes.
        """,
        agent=follow_up_partner,
        expected_output="Post-meeting activation kit featuring action tracker, follow-up messaging, and enablement guidance."
    )

    meeting_prep_crew = Crew(
        agents=[
            context_analyzer,
            industry_insights_generator,
            strategy_formulator,
            executive_briefing_creator,
            rehearsal_coach,
            follow_up_partner,
        ],
        tasks=[
            context_analysis_task,
            industry_analysis_task,
            strategy_development_task,
            executive_brief_task,
            rehearsal_simulation_task,
            follow_up_activation_task,
        ],
        verbose=True,
        process=Process.sequential
    )

    if st.button("Prepare Meeting"):
        if not company_name or not meeting_objective:
            st.warning("Please provide both a company name and meeting objective before preparing the meeting.")
        else:
            with st.spinner("AI agents are preparing your meeting..."):
                try:
                    result = meeting_prep_crew.kickoff()
                except Exception as exc:
                    st.error(f"Error preparing meeting: {exc}")
                    result = None

            if result:
                result_text = _result_to_markdown(result)
                st.markdown(result_text)
                try:
                    st.download_button(
                        label="Download meeting brief as .md",
                        data=result_text,
                        file_name=f"meeting_brief_{company_name}.md",
                        mime="text/markdown"
                    )
                except Exception:
                    pass

                # Generate PPTX deck from the brief
                if st.button("Generate Slides (PPTX)"):
                    pptx_bytes = _brief_to_pptx(result_text, title=f"Meeting Brief - {company_name}")
                    if pptx_bytes:
                        try:
                            st.download_button(
                                label="Download slides",
                                data=pptx_bytes,
                                file_name=f"meeting_brief_{company_name}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key="download_pptx"
                            )
                        except Exception:
                            pass
                    else:
                        st.warning("python-pptx not available. Please install dependencies to enable slide generation.")

                timestamp = datetime.utcnow().isoformat()
                history_entry = {
                    "timestamp": timestamp,
                    "company": company_name,
                    "objective": meeting_objective,
                    "attendees": attendees,
                    "focusAreas": focus_areas,
                    "documents": [doc["name"] for doc in supporting_documents],
                    "result": result_text,
                }
                _save_meeting_to_history(history_entry)
                st.session_state["history_view"] = history_entry
                st.success("Meeting brief archived in your library.")

    # Practice Mode (text-only)
    with st.expander("Practice Mode (text-only)"):
        if "practice_history" not in st.session_state:
            st.session_state["practice_history"] = []  # list of {role: 'coach'|'you', content: str}

        # Show conversation
        if st.session_state["practice_history"]:
            st.markdown("### Session log")
            for turn in st.session_state["practice_history"]:
                role = turn.get("role", "coach")
                content = turn.get("content", "")
                prefix = "Coach" if role == "coach" else "You"
                st.markdown(f"**{prefix}:**\n\n{content}")

        user_resp = st.text_area("Your response", key="practice_user_response", height=120)
        c1, c2, c3 = st.columns(3)
        with c1:
            ask_next = st.button("Ask next objection", key="ask_next_objection")
        with c2:
            score_it = st.button("Score my response", key="score_response")
        with c3:
            clear_sess = st.button("Clear session", key="clear_practice_session")

        if clear_sess:
            st.session_state["practice_history"] = []
            st.experimental_rerun()

        # Generate next objection
        if ask_next:
            history_text = "\n".join(
                [f"{t.get('role', 'coach')}: {t.get('content', '')}" for t in st.session_state["practice_history"]][-10:]
            )
            task = Task(
                description=f"""
                Generate the next realistic stakeholder objection for a rehearsal.
                Company: {company_name}
                Objective: {meeting_objective}
                Attendees: {attendees}
                Focus areas: {focus_areas}
                Recent practice log (last 10 turns):
                {history_text}

                Output 1-2 sentences with a sharp, persona-driven objection.
                """,
                agent=rehearsal_coach,
                expected_output="A concise, realistic objection in 1-2 sentences."
            )
            temp_crew = Crew(agents=[rehearsal_coach], tasks=[task], verbose=False, process=Process.sequential)
            with st.spinner("Thinking of a tough objection..."):
                try:
                    obj_res = temp_crew.kickoff()
                except Exception as exc:
                    st.error(f"Error generating objection: {exc}")
                    obj_res = None
            if obj_res:
                objection = _result_to_markdown(obj_res)
                st.session_state["practice_history"].append({"role": "coach", "content": objection})
                st.experimental_rerun()

        # Score user's response
        if score_it and user_resp.strip():
            st.session_state["practice_history"].append({"role": "you", "content": user_resp})
            history_text = "\n".join(
                [f"{t.get('role', 'coach')}: {t.get('content', '')}" for t in st.session_state["practice_history"]][-10:]
            )
            task = Task(
                description=f"""
                Evaluate the user's response to the last objection.
                Provide:
                - Score (1-10) on clarity, evidence, and relevance
                - 3 coaching tips to improve
                - A refined sample answer (3-5 sentences)

                Context: {company_name}, objective: {meeting_objective}, focus areas: {focus_areas}
                Practice log (last 10 turns):
                {history_text}
                User response:
                {user_resp}
                """,
                agent=rehearsal_coach,
                expected_output="Markdown with a short rubric, 3 tips, and a refined sample answer."
            )
            temp_crew = Crew(agents=[rehearsal_coach], tasks=[task], verbose=False, process=Process.sequential)
            with st.spinner("Scoring your response..."):
                try:
                    score_res = temp_crew.kickoff()
                except Exception as exc:
                    st.error(f"Error scoring response: {exc}")
                    score_res = None
            if score_res:
                feedback = _result_to_markdown(score_res)
                st.session_state["practice_history"].append({"role": "coach", "content": feedback})
                st.experimental_rerun()

    history_preview = st.session_state.get("history_view")
    if history_preview:
        with st.expander("Saved brief preview", expanded=False):
            st.markdown(history_preview["result"])
            try:
                st.download_button(
                    label="Download selected brief",
                    data=history_preview["result"],
                    file_name=f"saved_meeting_brief_{history_preview['company']}.md",
                    mime="text/markdown",
                    key="download_history_brief"
                )
            except Exception:
                pass

    st.sidebar.markdown("""
    ## How to use this app:
    1. Enter your API keys in the sidebar.
    2. Provide details about the meeting and optional supporting assets.
    3. Click 'Prepare Meeting' to generate the preparation package.

    New premium capabilities:
    - Store and replay briefs from the meeting library.
    - Fuse uploaded documents and historical notes into every analysis.
    - Run a rehearsal simulation with objection handling guidance.
    - Produce post-meeting activation plans and follow-up messaging.
    """)
else:
    st.warning("Please enter your OpenAI and Serper API keys in the sidebar before proceeding.")